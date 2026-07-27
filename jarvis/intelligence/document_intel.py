"""Document intelligence extensions — more formats, tagging, entity extraction."""

from __future__ import annotations

import csv
import logging
import re
from pathlib import Path
from typing import Any

log = logging.getLogger("jarvis.intelligence.document_intel")

EXTENDED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".markdown",
    ".html",
    ".htm",
    ".csv",
    ".xlsx",
    ".pptx",
    ".py",
    ".js",
    ".ts",
    ".json",
    ".yml",
    ".yaml",
    ".xml",
    ".rst",
}


def supported_extensions() -> list[str]:
    return sorted(EXTENDED_EXTENSIONS)


def _clean(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", (text or "").replace("\r\n", "\n")).strip()


def extract_csv(path: Path, *, max_rows: int = 200) -> dict[str, Any]:
    rows: list[list[str]] = []
    with path.open(newline="", encoding="utf-8", errors="ignore") as fh:
        reader = csv.reader(fh)
        for i, row in enumerate(reader):
            rows.append([str(c) for c in row])
            if i >= max_rows:
                break
    header = rows[0] if rows else []
    preview = "\n".join(", ".join(r) for r in rows[:40])
    return {
        "ok": True,
        "title": path.stem,
        "source": "csv",
        "pages": [preview],
        "tables": [{"headers": header, "row_count": max(0, len(rows) - 1)}],
        "char_count": len(preview),
    }


def extract_xlsx(path: Path, *, max_sheets: int = 5, max_rows: int = 100) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        return {"ok": False, "error": f"openpyxl required: {exc}"}
    wb = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[str] = []
    tables: list[dict[str, Any]] = []
    for sheet in wb.worksheets[:max_sheets]:
        rows = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                rows.append(vals)
            if i >= max_rows:
                break
        text = f"# Sheet: {sheet.title}\n" + "\n".join(" | ".join(r) for r in rows[:50])
        pages.append(_clean(text))
        tables.append({"sheet": sheet.title, "row_count": len(rows)})
    full = "\n\n".join(pages)
    return {
        "ok": True,
        "title": path.stem,
        "source": "openpyxl",
        "pages": pages or [""],
        "tables": tables,
        "char_count": len(full),
    }


def extract_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        return {"ok": False, "error": f"python-pptx required: {exc}"}
    prs = Presentation(str(path))
    pages: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        pages.append(_clean(f"[Slide {i}]\n" + "\n".join(parts)))
    full = "\n\n".join(pages)
    return {
        "ok": True,
        "title": path.stem,
        "source": "python-pptx",
        "pages": pages or [""],
        "char_count": len(full),
        "page_count": len(pages),
    }


def extract_code_or_text(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return {
        "ok": True,
        "title": path.name,
        "source": "text",
        "pages": [_clean(text)],
        "char_count": len(text),
    }


def parse_extended(path: str | Path) -> dict[str, Any]:
    p = Path(path).expanduser().resolve()
    if not p.is_file():
        return {"ok": False, "error": f"not found: {p}"}
    suffix = p.suffix.lower()
    if suffix not in EXTENDED_EXTENSIONS:
        return {"ok": False, "error": f"unsupported: {suffix}", "supported": supported_extensions()}

    # Prefer core pipeline for classic types
    if suffix in {".pdf", ".docx", ".txt", ".md", ".markdown", ".html", ".htm"}:
        try:
            from jarvis.document_pipeline import parse_document

            doc = parse_document(p)
            return {
                "ok": True,
                "title": doc.title,
                "source": doc.source,
                "pages": doc.pages,
                "char_count": doc.char_count,
                "page_count": doc.page_count,
                "path": str(p),
            }
        except Exception as exc:
            return {"ok": False, "error": str(exc)}

    if suffix == ".csv":
        result = extract_csv(p)
    elif suffix == ".xlsx":
        result = extract_xlsx(p)
    elif suffix == ".pptx":
        result = extract_pptx(p)
    else:
        result = extract_code_or_text(p)
    if result.get("ok"):
        result["path"] = str(p)
        result["page_count"] = result.get("page_count") or len(result.get("pages") or [])
    return result


def auto_tags(text: str, *, limit: int = 12) -> list[str]:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_-]{3,}", text or "")
    freq: dict[str, int] = {}
    stop = {
        "this",
        "that",
        "with",
        "from",
        "have",
        "were",
        "been",
        "will",
        "your",
        "about",
        "into",
        "they",
        "them",
        "then",
        "than",
        "also",
        "only",
        "just",
        "like",
        "aria",
    }
    for w in words:
        key = w.lower()
        if key in stop:
            continue
        freq[key] = freq.get(key, 0) + 1
    ranked = sorted(freq.items(), key=lambda x: (-x[1], x[0]))
    return [w for w, _ in ranked[:limit]]


def analyze_document(path: str | Path) -> dict[str, Any]:
    parsed = parse_extended(path)
    if not parsed.get("ok"):
        return parsed
    text = "\n\n".join(parsed.get("pages") or [])
    from jarvis.intelligence.knowledge_graph import extract_entities, extract_relationships

    tags = auto_tags(text)
    entities = extract_entities(text)
    relationships = extract_relationships(text)
    summary = text[:600].rstrip() + ("…" if len(text) > 600 else "")
    return {
        **parsed,
        "tags": tags,
        "entities": entities,
        "relationships": relationships,
        "summary": summary,
    }


def ocr_image(path: str | Path) -> dict[str, Any]:
    """OCR via vision module when available; degrade gracefully."""
    p = Path(path).expanduser()
    if not p.is_file():
        return {"ok": False, "error": f"not found: {p}"}
    try:
        from jarvis.modules import vision

        fn = getattr(vision, "ocr", None) or getattr(vision, "ocr_image", None) or getattr(vision, "describe", None)
        if not callable(fn):
            return {"ok": False, "error": "vision OCR unavailable"}
        result = fn(str(p))
        if isinstance(result, dict):
            return {"ok": True, **result}
        return {"ok": True, "text": str(result)}
    except Exception as exc:
        log.warning("ocr failed: %s", exc)
        return {"ok": False, "error": str(exc), "hint": "Enable vision model or attach as chat OCR request"}
